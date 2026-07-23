import os
import re
import json
from datetime import datetime
from typing import List, Dict, Any, Tuple, Optional

# IMPORTS OPCIONALES Y SEGUROS PARA MÚLTIPLES FORMATOS DE DOCUMENTOS
try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

# OCR Opcional para PDFs escaneados
try:
    import pytesseract
    from pdf2image import convert_from_path
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False


from app.services.category_registry import resolve_folder_metadata


class DocumentChunk:
    """
    =============================================================================
    ESTRUCTURA DE UN FRAGMENTO DE TEXTO Y SUS METADATOS COMPLETOS (ETAPA 2)
    =============================================================================
    Contiene el contenido extraído y los metadatos asignados para Gobernanza RAG:
    - Categoría / Área de negocio
    - Responsable oficial (Ownership)
    - Ubicación exacta (Página, Diapositiva, Hoja/Fila)
    - Nombre del archivo y Fecha de Modificación
    """
    def __init__(
        self,
        content: str,
        source: str,
        category: str,
        owner: str,
        chunk_id: int,
        location: str = "General",
        modified_at: str = "",
        author: str = "",
    ):
        self.content = content         # Texto del fragmento limpio
        self.source = source           # Nombre del archivo fuente (ej: manual_rh.pdf)
        self.category = category       # Categoría/Área (ej: 'Recursos Humanos', 'Finanzas')
        self.owner = owner             # Responsable oficial / Ownership (ej: rh@aluratech.com)
        self.chunk_id = chunk_id       # ID secuencial del fragmento dentro del archivo
        self.location = location       # Ubicación exacta (ej: 'Página 3', 'Diapositiva 2')
        self.modified_at = modified_at # Fecha de modificación del archivo (YYYY-MM-DD)
        self.author = author             # Autor del documento cuando está disponible en metadatos

    def to_dict(self) -> Dict[str, Any]:
        return {
            "content": self.content,
            "source": self.source,
            "category": self.category,
            "owner": self.owner,
            "chunk_id": self.chunk_id,
            "location": self.location,
            "modified_at": self.modified_at,
            "author": self.author,
        }

    @staticmethod
    def from_dict(data: Dict[str, Any]) -> "DocumentChunk":
        return DocumentChunk(
            content=data["content"],
            source=data["source"],
            category=data["category"],
            owner=data.get("owner", "soporte@aluratech.com"),
            chunk_id=int(data.get("chunk_id", 0)),
            location=data.get("location", "General"),
            modified_at=data.get("modified_at", ""),
            author=data.get("author", ""),
        )


class DocumentLoader:
    """
    =============================================================================
    PROCESAMIENTO Y EXTRACCIÓN DE CONTENIDO (ETAPA 2)
    =============================================================================
    Fase encargada de transformar los documentos originales en sus variados formatos
    en texto limpio y estructurado con sus respectivos metadatos de gobernanza:
    
    1. Extracción por formato: PDF (nativo + OCR), Word, Excel, PowerPoint (+ notas del orador),
       Markdown, CSV, JSON y HTML.
    2. Limpieza del texto: Eliminación de ruidos, encabezados/pies repetidos, caracteres especiales
       y espacios corrompidos.
    3. Chunking inteligente: División en bloques con overlap manteniendo contexto semántico.
    4. Atribución de metadatos: Asignación de categoría, ownership, ubicación exacta y fechas.
    """

    @staticmethod
    def infer_metadata_from_path(file_path: str) -> Tuple[str, str]:
        """Deduce Categoría y Responsable oficial según la subcarpeta fuente."""
        parent_dir = os.path.basename(os.path.dirname(file_path)).lower()
        return resolve_folder_metadata(parent_dir)

    @staticmethod
    def extract_document_author(file_path: str) -> str:
        """Lee autor desde metadatos del archivo cuando el formato lo expone (Word)."""
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".docx" and DocxDocument:
            try:
                doc = DocxDocument(file_path)
                author = (doc.core_properties.author or "").strip()
                return author
            except Exception:
                return ""
        return ""

    @staticmethod
    def get_file_modified_date(file_path: str) -> str:
        """Obtiene la fecha de última modificación del archivo en formato YYYY-MM-DD."""
        try:
            mtime = os.path.getmtime(file_path)
            return datetime.fromtimestamp(mtime).strftime("%Y-%m-%d")
        except Exception:
            return datetime.now().strftime("%Y-%m-%d")

    @staticmethod
    def clean_text(text: str) -> str:
        """
        FASE 2: LIMPIEZA DE TEXTO
        Elimina ruido que no aporta significado:
        - Caracteres especiales de formato corrupto (\x00, control chars)
        - Encabezados/pies de página tipo 'Página X de Y' o 'Confidencial'
        - Espacios múltiples duplicados y saltos de línea excesivos
        """
        if not text:
            return ""

        # Eliminar caracteres nulos o no imprimibles
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)

        # Eliminar patrones comunes de números de página y encabezados ruidosos
        text = re.sub(r'(?i)p[áa]gina\s+\d+(\s+de\s+\d+)?', '', text)
        text = re.sub(r'(?i)confidencial\s+-\s+uso\s+interno', '', text)

        # Normalizar espacios tabulares y múltiples a un solo espacio
        text = re.sub(r'[ \t]+', ' ', text)

        # Normalizar saltos de línea triples o mayores a doble salto de línea
        text = re.sub(r'\n{3,}', '\n\n', text)

        return text.strip()

    # =========================================================================
    # FASE 1: EXTRACCIÓN POR FORMATO ESTRUCTURADA POR UBICACIÓN
    # =========================================================================

    @staticmethod
    def load_and_structure_file(file_path: str) -> List[Dict[str, str]]:
        """
        Extrae el contenido formateado y estructurado con su ubicación exacta dentro del archivo.
        Retorna lista de dicts: [{"content": "...", "location": "Página 1"}, ...]
        """
        filename = os.path.basename(file_path)
        ext = os.path.splitext(filename)[1].lower()

        if ext == ".pdf":
            return DocumentLoader._read_pdf_structured(file_path)
        elif ext == ".docx":
            return DocumentLoader._read_docx_structured(file_path)
        elif ext == ".pptx":
            return DocumentLoader._read_pptx_structured(file_path)
        elif ext in [".xlsx", ".xls"]:
            return DocumentLoader._read_excel_structured(file_path)
        elif ext == ".csv":
            return DocumentLoader._read_csv_structured(file_path)
        elif ext == ".json":
            return DocumentLoader._read_json_structured(file_path)
        elif ext in [".html", ".htm"]:
            return DocumentLoader._read_html_structured(file_path)
        elif ext in [".md", ".txt"]:
            return DocumentLoader._read_markdown_structured(file_path)
        else:
            print(f"[DocumentLoader] ⚠️ Formato omitido: '{ext}' en {filename}")
            return []

    @staticmethod
    def _read_pdf_structured(file_path: str) -> List[Dict[str, str]]:
        """PDF: Extracción nativa página a página con fallback OCR para PDF escaneados."""
        blocks = []
        if not PdfReader:
            return blocks

        try:
            reader = PdfReader(file_path)
            for page_idx, page in enumerate(reader.pages, start=1):
                raw_text = page.extract_text() or ""
                cleaned = DocumentLoader.clean_text(raw_text)

                # Fallback OCR si la página no contiene texto reconocible y OCR está disponible
                if len(cleaned) < 15 and OCR_AVAILABLE:
                    try:
                        images = convert_from_path(file_path, first_page=page_idx, last_page=page_idx)
                        if images:
                            ocr_text = pytesseract.image_to_string(images[0], lang='spa+eng')
                            cleaned = DocumentLoader.clean_text(ocr_text)
                    except Exception as ocr_err:
                        print(f"[DocumentLoader] ⚠️ Error en OCR para PDF {file_path} pág {page_idx}: {ocr_err}")

                if cleaned:
                    blocks.append({"content": cleaned, "location": f"Página {page_idx}"})
        except Exception as e:
            print(f"[DocumentLoader] ❌ Error extrayendo PDF {file_path}: {e}")

        return blocks

    @staticmethod
    def _read_docx_structured(file_path: str) -> List[Dict[str, str]]:
        """WORD: Extracción de párrafos preservando títulos y tablas estructuradas."""
        blocks = []
        if not DocxDocument:
            return blocks

        try:
            doc = DocxDocument(file_path)
            current_section = "Sección General"
            section_text = []

            for paragraph in doc.paragraphs:
                p_text = paragraph.text.strip()
                if not p_text:
                    continue

                # Detectar estilos de títulos/encabezados para cambiar de sección
                if paragraph.style.name.startswith("Heading") or paragraph.style.name.startswith("Título"):
                    if section_text:
                        full_content = DocumentLoader.clean_text("\n".join(section_text))
                        if full_content:
                            blocks.append({"content": full_content, "location": current_section})
                        section_text = []
                    current_section = f"Sección: {p_text}"
                else:
                    section_text.append(p_text)

            if section_text:
                full_content = DocumentLoader.clean_text("\n".join(section_text))
                if full_content:
                    blocks.append({"content": full_content, "location": current_section})

            # Extraer tablas incorporadas en el documento Word
            for t_idx, table in enumerate(doc.tables, start=1):
                table_rows = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_data:
                        table_rows.append(" | ".join(row_data))
                if table_rows:
                    blocks.append({
                        "content": DocumentLoader.clean_text("\n".join(table_rows)),
                        "location": f"Tabla {t_idx} ({current_section})"
                    })
        except Exception as e:
            print(f"[DocumentLoader] ❌ Error extrayendo DOCX {file_path}: {e}")

        return blocks

    @staticmethod
    def _read_pptx_structured(file_path: str) -> List[Dict[str, str]]:
        """POWERPOINT: Extracción de texto por diapositiva + notas del orador (Speaker Notes)."""
        blocks = []
        if not Presentation:
            return blocks

        try:
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_content = []

                # Extracción del contenido de formas y cajas de texto
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        slide_content.append(shape.text.strip())

                # Extracción crítica de las notas del orador
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_content.append(f"[Notas del Orador]: {notes}")

                cleaned = DocumentLoader.clean_text("\n".join(slide_content))
                if cleaned:
                    blocks.append({"content": cleaned, "location": f"Diapositiva {slide_idx}"})
        except Exception as e:
            print(f"[DocumentLoader] ❌ Error extrayendo PPTX {file_path}: {e}")

        return blocks

    @staticmethod
    def _read_excel_structured(file_path: str) -> List[Dict[str, str]]:
        """EXCEL: Convierte planillas en texto estructurado fila a fila con encabezados repetidos."""
        blocks = []
        try:
            if pd is not None:
                excel_file = pd.ExcelFile(file_path)
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    df = df.dropna(how="all")
                    rows_text = []
                    for row_idx, row in df.iterrows():
                        formatted_row = ", ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                        if formatted_row:
                            rows_text.append(f"[Fila {row_idx + 2}] {formatted_row}")
                    
                    cleaned = DocumentLoader.clean_text("\n".join(rows_text))
                    if cleaned:
                        blocks.append({"content": cleaned, "location": f"Hoja Excel '{sheet_name}'"})
        except Exception as e:
            print(f"[DocumentLoader] ❌ Error extrayendo Excel {file_path}: {e}")

        return blocks

    @staticmethod
    def _read_csv_structured(file_path: str) -> List[Dict[str, str]]:
        """CSV: Extrae filas asociando cada valor con su encabezado de columna."""
        blocks = []
        try:
            if pd is not None:
                df = pd.read_csv(file_path)
                rows_text = []
                for row_idx, row in df.iterrows():
                    formatted_row = ", ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    if formatted_row:
                        rows_text.append(f"[Fila {row_idx + 1}] {formatted_row}")
                cleaned = DocumentLoader.clean_text("\n".join(rows_text))
                if cleaned:
                    blocks.append({"content": cleaned, "location": "Tabla CSV"})
        except Exception as e:
            print(f"[DocumentLoader] ❌ Error extrayendo CSV {file_path}: {e}")

        return blocks

    @staticmethod
    def _read_json_structured(file_path: str) -> List[Dict[str, str]]:
        """JSON: Transforma estructuras jerárquicas en texto comprensible para el agente."""
        blocks = []
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            if isinstance(data, list):
                items = [json.dumps(item, ensure_ascii=False) for item in data]
                cleaned = DocumentLoader.clean_text("\n".join(items))
                blocks.append({"content": cleaned, "location": "Estructura JSON (Lista)"})
            elif isinstance(data, dict):
                formatted = "\n".join([f"{k}: {json.dumps(v, ensure_ascii=False)}" for k, v in data.items()])
                cleaned = DocumentLoader.clean_text(formatted)
                blocks.append({"content": cleaned, "location": "Objeto JSON"})
        except Exception as e:
            print(f"[DocumentLoader] ❌ Error extrayendo JSON {file_path}: {e}")

        return blocks

    @staticmethod
    def _read_html_structured(file_path: str) -> List[Dict[str, str]]:
        """HTML: Elimina marcas técnicas y etiquetas conservando el texto semántico."""
        blocks = []
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                raw_html = f.read()

            if BeautifulSoup:
                soup = BeautifulSoup(raw_html, "html.parser")
                for script in soup(["script", "style"]):
                    script.extract()
                text = soup.get_text(separator="\n", strip=True)
            else:
                text = re.sub(r'<[^>]+>', ' ', raw_html)

            cleaned = DocumentLoader.clean_text(text)
            if cleaned:
                blocks.append({"content": cleaned, "location": "Documento HTML"})
        except Exception as e:
            print(f"[DocumentLoader] ❌ Error extrayendo HTML {file_path}: {e}")

        return blocks

    @staticmethod
    def _read_markdown_structured(file_path: str) -> List[Dict[str, str]]:
        """MARKDOWN / TXT: Extrae texto dividiendo por títulos de Markdown (# ##)."""
        blocks = []
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()

            lines = content.splitlines()
            current_section = "Sección General"
            section_lines = []

            for line in lines:
                if line.strip().startswith("#"):
                    if section_lines:
                        cleaned = DocumentLoader.clean_text("\n".join(section_lines))
                        if cleaned:
                            blocks.append({"content": cleaned, "location": current_section})
                        section_lines = []
                    current_section = f"Sección: {line.strip('#').strip()}"
                else:
                    section_lines.append(line)

            if section_lines:
                cleaned = DocumentLoader.clean_text("\n".join(section_lines))
                if cleaned:
                    blocks.append({"content": cleaned, "location": current_section})

        except Exception as e:
            print(f"[DocumentLoader] ❌ Error extrayendo Markdown {file_path}: {e}")

        return blocks

    # =========================================================================
    # FASE 3 Y 4: CHUNKING CON SOBREPOSICIÓN Y ATRIBUCIÓN DE METADATOS
    # =========================================================================

    @staticmethod
    def _chunk_text_with_overlap(text: str, chunk_size: int, overlap: int) -> List[str]:
        """Divide texto por caracteres con solapamiento (overlap)."""
        if not text:
            return []
        if len(text) <= chunk_size:
            return [text]

        chunks: List[str] = []
        start = 0
        while start < len(text):
            end = min(start + chunk_size, len(text))
            chunks.append(text[start:end].strip())
            if end >= len(text):
                break
            start += max(1, chunk_size - overlap)
        return [c for c in chunks if c]

    @staticmethod
    def process_and_chunk_file(
        file_path: str,
        chunk_size: int = 1000,
        overlap: int = 150,
    ) -> List[DocumentChunk]:
        """
        PROCESAMIENTO COMPLETO DE INGESTA (FASE 1 A 4):
        1. Extracción por formato (PDF, Word, PPTX, Excel, CSV, JSON, HTML, MD)
        2. Limpieza de texto profunda
        3. Chunking con superposición (overlap) para no cortar ideas
        4. Atribución de metadatos (Categoría, Ownership, Ubicación, Fecha Modificación)
        """
        filename = os.path.basename(file_path)
        category, owner = DocumentLoader.infer_metadata_from_path(file_path)
        modified_at = DocumentLoader.get_file_modified_date(file_path)
        author = DocumentLoader.extract_document_author(file_path)

        structured_blocks = DocumentLoader.load_and_structure_file(file_path)
        if not structured_blocks:
            return []

        all_chunks = []
        global_chunk_id = 0

        for block in structured_blocks:
            text = block["content"]
            location = block["location"]

            for chunk_text in DocumentLoader._chunk_text_with_overlap(text, chunk_size, overlap):
                all_chunks.append(DocumentChunk(
                    content=chunk_text,
                    source=filename,
                    category=category,
                    owner=owner,
                    chunk_id=global_chunk_id,
                    location=location,
                    modified_at=modified_at,
                    author=author,
                ))
                global_chunk_id += 1

        return all_chunks
